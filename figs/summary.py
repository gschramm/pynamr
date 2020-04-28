import os
import pptx

from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

from glob import glob

sdirs    = ['inverse_crime','50','8']
f1_files = []

for sdir in sdirs:
  # add prior recons with no noise
  f1_files += sorted(glob(os.path.join(sdir,'*beta_3.0_*noise_level_0.0_*f1.png')))
  f1_files += sorted(glob(os.path.join(sdir,'*beta_10.0_*noise_level_0.0_*f1.png')))
  f1_files += sorted(glob(os.path.join(sdir,'*beta_30.0_*noise_level_0.0_*f1.png')))
  
  # add prior recons with noise
  f1_files += sorted(glob(os.path.join(sdir,'*beta_3.0_*noise_level_0.1_*f1.png')))
  f1_files += sorted(glob(os.path.join(sdir,'*beta_10.0_*noise_level_0.1_*f1.png')))
  f1_files += sorted(glob(os.path.join(sdir,'*beta_30.0_*noise_level_0.1_*f1.png')))

  f1_files += sorted(glob(os.path.join(sdir,'*beta_3.0_*noise_level_0.5_*f1.png')))
  f1_files += sorted(glob(os.path.join(sdir,'*beta_10.0_*noise_level_0.5_*f1.png')))
  f1_files += sorted(glob(os.path.join(sdir,'*beta_30.0_*noise_level_0.5_*f1.png')))

f2_files = [x.replace('f1.png','f2.png') for x in f1_files]

prs = pptx.Presentation('template.pptx')
prs.slide_width  = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(4.5)

for i, f1_file in enumerate(f1_files):
  T2star_recon_short = os.path.basename(f1_file).split('T2star_recon_short_')[1].split('__')[0]
  noise_level        = os.path.basename(f1_file).split('noise_level_')[1].split('__')[0]
  beta               = os.path.basename(f1_file).split('beta_')[1].split('__')[0]

  if T2star_recon_short == '-1.0':
    t2star_model = 'inverse crime'
  elif T2star_recon_short == '8.0':
    t2star_model = 'GM'
  elif T2star_recon_short == '50.0':
    t2star_model = 'CSF'
  else:
    t2star_model = 'unknown'
    print('unknown model')

  blank_slide_layout = prs.slide_layouts[8]
  slide = prs.slides.add_slide(blank_slide_layout)
 
  pic = slide.shapes.add_picture(f1_file, pptx.util.Inches(0), pptx.util.Inches(1.),            
                                 height = pptx.util.Inches(3))

  pic = slide.shapes.add_picture(f2_files[i], pptx.util.Inches(4), pptx.util.Inches(1.),            
                                 height = pptx.util.Inches(3))

  title_placeholder = slide.shapes.title
  title_placeholder.text = f'T2* model {t2star_model}   beta {beta}   noise level {noise_level}'
  title_placeholder.text_frame.paragraphs[0].font.size = Pt(14)
  title_placeholder.text_frame.paragraphs[0].font.bold = True

  #p = title_placeholder.text_frame.paragraphs[0]
  #run = p.add_run()
  #font = run.font
  #font.name = 'Calibri'
  #font.size = Pt(18)
  #font.bold = True


  #title_placeholder.margin_bottom = pptx.util.Inches(2)
  #title_placeholder.margin_left = 0
  #title_placeholder.vertical_anchor = MSO_ANCHOR.TOP
  #title_placeholder.word_wrap = False
  #title_placeholder.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

prs.save("summary.pptx")
